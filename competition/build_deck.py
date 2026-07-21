# -*- coding: utf-8 -*-
"""
尋形 Conform — Claude Code 開發競賽簡報 — 6 分鐘簡報 + Q&A 另計（約 3 分鐘）(16:9)
產出: 尋形Conform_競賽簡報.pptx（封面 + 8 內容頁 + Q&A + BACKUP 決策表，共 11 張）
風格: 參考 contrastive_learning_evolution 風格 —
      冷灰底 #F4F6FB + 深 slate 卡 #162039、膠囊徽章、Tailwind 重點色、
      PingFang TC 標題 + Menlo 等寬標籤、卡片/架構流程語言。
重點: 無頁數/章節限制; 以 6:00 節奏配頁, 每頁 notes 開頭標 ⏱ 時間帶;
      決策表併入主檔 BACKUP 頁（Q&A 免切檔）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ---------- palette (Tailwind on slate, lifted from reference deck) ----------
BG       = RGBColor(0xF4, 0xF6, 0xFB)   # light cool-gray page
BG_DARK  = RGBColor(0x0E, 0x16, 0x26)   # deep ink (dark cards / hero)
NAVY     = RGBColor(0x16, 0x20, 0x39)   # slate card
NAVY2    = RGBColor(0x23, 0x37, 0x5C)   # lighter slate
HEAD     = RGBColor(0x16, 0x21, 0x3A)   # heading ink on light
TEXT     = RGBColor(0x1F, 0x2D, 0x45)   # body on light
MUTE     = RGBColor(0x5C, 0x6B, 0x8A)   # caption / secondary
DLT      = RGBColor(0xC6, 0xD6, 0xF2)   # body on dark
DMUTE    = RGBColor(0x90, 0xA0, 0xC0)   # caption on dark
CARDLN   = RGBColor(0xDC, 0xE3, 0xF0)   # card border
TINT     = RGBColor(0xEA, 0xF0, 0xFA)   # light-blue tint card

AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)
CYAN6    = RGBColor(0x06, 0xB6, 0xD4)
CYAN8    = RGBColor(0x0E, 0x74, 0x90)
ROSE     = RGBColor(0xF4, 0x3F, 0x5E)
ROSE_DK  = RGBColor(0xC0, 0x39, 0x5A)
ROSE_LT  = RGBColor(0xFD, 0xE7, 0xEB)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
EMER     = RGBColor(0x2A, 0x8F, 0x5F)
EMER_LT  = RGBColor(0xE4, 0xF4, 0xEC)
AMBER_LT = RGBColor(0xFD, 0xF0, 0xDB)
GRAYBG   = RGBColor(0xEE, 0xF2, 0xF8)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "PingFang TC"     # CJK + Latin body/heading
MONO = "Menlo"           # labels, years, code, page numbers

TOTAL = 9                # numbered pages: 封面 01 … 收尾 09 (Q&A / BACKUP 不編號)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
MIN_PT = 12   # 競賽規範: 全簡報字體不得小於 12pt — 由此處全域強制

def _setfont(run, size, bold=False, color=TEXT, name=FONT, italic=False):
    f = run.font
    size = max(size, MIN_PT)
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        el = sp._element.spPr
        ns = ('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
              '<a:outerShdw blurRad="55000" dist="20000" dir="5400000" rotWithShape="0">'
              '<a:srgbClr val="2A3A55"><a:alpha val="26000"/></a:srgbClr></a:outerShdw></a:effectLst>')
        el.append(parse_xml(ns))
    return sp


def txt(s, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', align)
        if ln.get('space_after') is not None:
            p.space_after = Pt(ln['space_after'])
        if ln.get('space_before') is not None:
            p.space_before = Pt(ln['space_before'])
        if ln.get('line_spacing') is not None:
            p.line_spacing = ln['line_spacing']
        for r in ln['runs']:
            run = p.add_run(); run.text = r[0]
            _setfont(run, r[1], r[2] if len(r) > 2 else False,
                     r[3] if len(r) > 3 else TEXT, name=r[5] if len(r) > 5 else FONT,
                     italic=r[4] if len(r) > 4 else False)
    return tb


def one(text, size, bold=False, color=TEXT, italic=False, align=PP_ALIGN.LEFT, sa=None, sb=None, ls=None, name=FONT):
    return {'runs': [(text, size, bold, color, italic, name)], 'align': align, 'space_after': sa,
            'space_before': sb, 'line_spacing': ls}


def card(s, x, y, w, h, kind='white', shadow=True, radius=0.05):
    if kind == 'white':
        rect(s, x, y, w, h, fill=WHITE, line=CARDLN, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=shadow, radius=radius)
    elif kind == 'navy':
        rect(s, x, y, w, h, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=shadow, radius=radius)
    elif kind == 'dark':
        rect(s, x, y, w, h, fill=BG_DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=shadow, radius=radius)
    elif kind == 'amber':
        rect(s, x, y, w, h, fill=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=shadow, radius=radius)
    elif kind == 'tint':
        rect(s, x, y, w, h, fill=TINT, line=CARDLN, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=shadow, radius=radius)


def pill(s, x, y, w, h, zh, en, fill, fg=WHITE):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5, shadow=True)
    runs = []
    if zh:
        runs.append((zh, 13, True, fg, False, FONT))
    if en:
        runs.append(("  " + en, 11, True, fg, False, MONO))
    txt(s, x, y - 0.01, w, h, [{'runs': runs, 'align': PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)


def brandmark(s, x, y, size, form_color=CYAN6, bracket_color=NAVY):
    """尋形 mark — frame-select marquee corners bracketing a 'form' tile (= 框選一個形)."""
    t = size * 0.085
    arm = size * 0.32
    fs = size * 0.46
    fo = (size - fs) / 2.0
    rect(s, x + fo, y + fo, fs, fs, fill=form_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.2)
    rect(s, x + size * 0.06, y + size * 0.06, arm, t, fill=bracket_color)
    rect(s, x + size * 0.06, y + size * 0.06, t, arm, fill=bracket_color)
    rect(s, x + size * 0.94 - arm, y + size * 0.94 - t, arm, t, fill=bracket_color)
    rect(s, x + size * 0.94 - t, y + size * 0.94 - arm, t, arm, fill=bracket_color)


def header(s, badge_zh, badge_en, zh, en, idx, accent=AMBER, badge_w=2.5, tag=None):
    rect(s, 0, 0, SW, SH, fill=BG)
    pill(s, 0.6, 0.5, badge_w, 0.46, badge_zh, badge_en, accent, WHITE if accent != AMBER else NAVY)
    txt(s, 0.57, 1.04, 9.4, 0.7, [one(zh, 27, True, HEAD)])
    txt(s, 0.6, 1.66, 9.6, 0.4, [one(en, 13, False, MUTE, italic=True)])
    # top-right brand lockup
    brandmark(s, SW - 0.92, 0.52, 0.36)
    txt(s, SW - 4.3, 0.46, 3.27, 0.34, [one("尋形 Conform", 13, True, HEAD, align=PP_ALIGN.RIGHT)], anchor=MSO_ANCHOR.MIDDLE)
    page = tag if tag else f"{idx:02d} / {TOTAL:02d}"
    txt(s, SW - 4.3, 0.78, 3.27, 0.3, [one(page, 11, True, accent, align=PP_ALIGN.RIGHT, name=MONO)], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 0.6, 2.0, SW - 1.2, 0.014, fill=CARDLN)


def label(s, x, y, text, color=CYAN6, size=11):
    """Menlo letterspaced eyebrow label (e.g. LOSS, LIVE DEMO)."""
    txt(s, x, y, 6, 0.3, [one(" ".join(text), size, True, color, name=MONO)])


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def chip(s, x, y, w, h, text, fill, fg, size=12, bold=True):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    txt(s, x + 0.05, y, w - 0.1, h, [one(text, size, bold, fg, align=PP_ALIGN.CENTER)],
        anchor=MSO_ANCHOR.MIDDLE)


# =================================================================
# Slide 1 — 封面 / Hook
# =================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=BG_DARK)
rect(s, 0, SH - 0.16, SW, 0.16, fill=CYAN6)
brandmark(s, 5.92, 1.15, 1.5, form_color=CYAN, bracket_color=WHITE)
txt(s, 0.6, 3.0, SW - 1.2, 1.0, [
    one("尋形 Conform", 48, True, WHITE, align=PP_ALIGN.CENTER)])
txt(s, 0.6, 4.05, SW - 1.2, 0.5, [
    one("框選一個形，找出所有同形——先進封裝 DXF 的幾何比對與設計規則檢查", 17, True, DLT, align=PP_ALIGN.CENTER)])
txt(s, 0.6, 4.62, SW - 1.2, 0.4, [
    one("Frame-select one pattern, find every instance — geometry matching & DRC for advanced packaging",
        12, False, DMUTE, align=PP_ALIGN.CENTER, italic=True)])
# teaser metrics strip
tw = 3.3; tg = 0.5
tx0 = (SW - 3 * tw - 2 * tg) / 2
teasers = [("480 → 1", "等效人工工時 → 機器 1 小時"),
           ("100%", "選定圖樣全 entity 覆蓋（vs 抽樣）"),
           ("820,212", "單張真實 DXF 實體數，秒級互動")]
for i, (big, sub) in enumerate(teasers):
    x = tx0 + i * (tw + tg)
    rect(s, x, 5.35, tw, 1.05, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12, shadow=True)
    txt(s, x, 5.5, tw, 0.5, [one(big, 22, True, CYAN, align=PP_ALIGN.CENTER, name=MONO)])
    txt(s, x + 0.15, 6.0, tw - 0.3, 0.35, [one(sub, 10.5, False, DMUTE, align=PP_ALIGN.CENTER)])
txt(s, SW - 3.5, 0.35, 3.0, 0.3, [one(f"01 / {TOTAL:02d}", 11, True, CYAN, align=PP_ALIGN.RIGHT, name=MONO)])
notes(s, "⏱ 0:00–0:18（18 秒）\n"
         "開場一句自我介紹＋一句產品定義：『尋形 Conform——在先進封裝的 DXF 圖紙上，框選一個圖樣，"
         "系統找出所有同形實例並自動跑設計規則檢查。』\n"
         "指一下三個數字當 hook（480→1、100% 覆蓋、82 萬實體秒級），細節後面每頁展開，這頁不停留。")

# =================================================================
# Slide 2 — 背景與痛點
# =================================================================
s = slide()
header(s, "痛點", "PROBLEM", "驗證瓶頸：先進封裝設計檢查的痛點",
       "The Verification Bottleneck — why Conform exists", 2, accent=ROSE, badge_w=2.45)

txt(s, 0.6, 2.16, 12.1, 0.5, [
    one("先進封裝後段驗證：一個漏檢的設計錯誤，可能造成基板重投、封裝組裝良率損失，甚至流到客戶端釀成客退——良率與可靠度的最後一道防線。", 14.5, True, TEXT)])

cards = [
    ("人力與時間", "6 位資深工程師", "× 超過 10 天", "最貴的人力，綁在最低槓桿的重複量測上"),
    ("覆蓋率", "抽樣式驗證", "缺陷會漏", "抽樣看不完 → 未覆蓋區就是漏檢風險"),
    ("授權瓶頸", "僅 4 套 AutoCAD", "併發授權", "授權數＝驗證並行度上限，全隊迭代卡死"),
]
cw = (12.13 - 2 * 0.34) / 3
cy = 2.72
for i, (head_, big, big2, sub) in enumerate(cards):
    x = 0.6 + i * (cw + 0.34)
    card(s, x, cy, cw, 2.78, 'white')
    rect(s, x, cy + 0.18, 0.12, 2.42, fill=ROSE)
    txt(s, x + 0.34, cy + 0.28, cw - 0.5, 0.5, [one(head_, 14, True, ROSE_DK)])
    txt(s, x + 0.34, cy + 0.92, cw - 0.5, 1.3, [
        one(big, 22, True, HEAD, sa=3),
        one(big2, 18, True, HEAD),
    ])
    txt(s, x + 0.34, cy + 2.12, cw - 0.5, 0.55, [one(sub, 13, False, MUTE, ls=1.05)])

card(s, 0.6, 5.78, 12.13, 0.92, 'navy')
txt(s, 0.95, 5.78, 11.4, 0.92, [
    one("這三個瓶頸是結構性的，不是效率問題——加人、加班、買授權只能緩解、無法消除，得換驗證方式。", 16, True, WHITE)],
    anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 0:18–0:58（40 秒）\n"
         "開場一句把嚴重性與『物理限制』框定：先進封裝漏檢後果嚴重；舊流程的三個限制不是『慢』，是結構性的——"
         "人力(6人>10天)、抽樣(看不完→會漏)、授權(僅4套AutoCAD併發，全隊序列化)。"
         "收尾關鍵句：這三個瓶頸只靠加人或加班解不開，必須換做法。下一頁帶出 尋形 Conform 的解。")

# =================================================================
# Slide 3 — 解法：框選到報告的閉環
# =================================================================
s = slide()
header(s, "解法", "SOLUTION", "解法：框選即查詢——從框選到規則報告的閉環",
       "Frame-select → match → rule check, one closed loop", 3, accent=CYAN6, badge_w=2.45)

txt(s, 0.6, 2.16, 12.1, 0.5, [
    one("工程師在瀏覽器裡框選一個圖樣，系統即時找出全部同形、確認後入庫；範本庫隨使用累積，之後新圖一鍵全掃、直出報告。",
        14.5, True, TEXT)])

flow3 = [
    ("① 上傳 DXF", "前處理＋圖層過濾＋渲染", NAVY, WHITE),
    ("② 框選圖樣", "AutoCAD 式互動操作", CYAN6, NAVY),
    ("③ 即時比對", "找出全部同形實例", NAVY, WHITE),
    ("④ 確認入庫", "範本庫隨使用累積", CYAN6, NAVY),
    ("⑤ 規則檢查", "100+ 條自動判 pass/fail", NAVY, WHITE),
]
fw = (12.13 - 4 * 0.4) / 5
fy = 2.78
for i, (t, sub, fill, fg) in enumerate(flow3):
    x = 0.6 + i * (fw + 0.4)
    rect(s, x, fy, fw, 1.3, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, shadow=True)
    sub_c = RGBColor(0x0B, 0x3A, 0x44) if fg == NAVY else DLT
    txt(s, x + 0.1, fy + 0.2, fw - 0.2, 0.5, [one(t, 14.5, True, fg, align=PP_ALIGN.CENTER)])
    txt(s, x + 0.12, fy + 0.72, fw - 0.24, 0.5, [one(sub, 11, False, sub_c, align=PP_ALIGN.CENTER, ls=1.05)])
    if i < 4:
        rect(s, x + fw + 0.08, fy + 0.52, 0.24, 0.26, fill=CYAN6, shape=MSO_SHAPE.CHEVRON)

# two supporting cards
sy3 = 4.4; sh3 = 1.8; sw3 = 5.95
card(s, 0.6, sy3, sw3, sh3, 'white')
rect(s, 0.6, sy3 + 0.18, 0.12, sh3 - 0.36, fill=CYAN6)
txt(s, 0.94, sy3 + 0.2, sw3 - 0.55, 0.4, [
    {'runs': [("零學習成本　", 14, True, CYAN8, False, FONT), ("FAMILIAR UX", 10, True, CYAN6, False, MONO)]}])
txt(s, 0.94, sy3 + 0.68, sw3 - 0.6, 1.0, [
    one("互動完全比照 AutoCAD：中鍵拖曳平移、左→右 window ／ 右→左 crossing 框選。", 12.5, False, TEXT, sa=4, ls=1.1),
    one("驗證工程師不需重新學一套工具，開瀏覽器就能上手。", 12.5, False, MUTE, ls=1.1)])

rx3 = 0.6 + sw3 + 0.33
card(s, rx3, sy3, sw3, sh3, 'dark')
label(s, rx3 + 0.32, sy3 + 0.2, "FLYWHEEL", AMBER, size=10)
txt(s, rx3 + 0.32, sy3 + 0.48, sw3 - 0.6, 0.4, [one("範本庫飛輪：越用越省", 14, True, WHITE)])
txt(s, rx3 + 0.32, sy3 + 0.94, sw3 - 0.64, 0.8, [
    one("每次確認的圖樣都累積進範本庫（現有 50+ 範本）；", 12, False, DLT, sa=4, ls=1.1),
    one("新圖上傳後一鍵 scan-all 全範本掃描——第一張圖教過的，之後每張圖都自動會。", 12, False, DLT, ls=1.1)])

card(s, 0.6, 6.42, 12.13, 0.72, 'navy')
txt(s, 0.95, 6.42, 11.4, 0.72, [
    one("閉環的關鍵在第③步「即時比對」——在 82 萬個實體裡找出全部同形，還要秒級回應。下一頁講怎麼辦到。",
        14.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 0:58–1:38（40 秒）\n"
         "用一個 operator 的一天講工作流程：上傳 DXF → 系統前處理渲染 → 框選一個圖樣 → 即時看到全部同形被標記 → "
         "確認入庫 → 規則引擎跑 100+ 條檢查直出報告。\n"
         "兩個賣點各一句：(1) 零學習成本——互動完全仿 AutoCAD（中鍵平移、window/crossing 框選），驗證工程師零遷移成本；"
         "(2) 範本庫飛輪——確認過的圖樣進庫（現有 50+ 範本），新圖一鍵 scan-all，教一次、之後全自動。\n"
         "收尾一句設 hook：閉環成立的前提是第③步能在 82 萬實體裡秒級找同形——轉下一頁演算法。")

# =================================================================
# Slide 4 — 演算法：規模與天文學靈感
# =================================================================
s = slide()
header(s, "演算法", "ALGORITHM", "像找星系一樣，在圖海中找出每一個同形",
       "The Matching Engine & its astronomy inspiration", 4, accent=CYAN6, badge_w=2.45)

# intro — bridge from the loop: step ③, and the hard part = scale
txt(s, 0.6, 2.14, 12.1, 0.46, [
    {'runs': [("閉環的第③步：在整張 DXF 裡找出框選圖樣的", 14.5, True, TEXT, False, FONT),
              ("全部同形實例", 14.5, True, CYAN8, False, FONT),
              ("。難點只有一個：", 14.5, True, TEXT, False, FONT),
              ("規模", 14.5, True, ROSE_DK, False, FONT),
              ("。", 14.5, True, TEXT, False, FONT)]}])

# two cards: brute-force wall (left) + astronomy idea (right)
ay = 2.6; ah = 2.3; aw = 5.95
card(s, 0.6, ay, aw, ah, 'white')
rect(s, 0.6, ay + 0.2, 0.12, ah - 0.4, fill=ROSE)
txt(s, 0.94, ay + 0.22, aw - 0.5, 0.4, [
    {'runs': [("暴力解的代價　", 14.5, True, ROSE_DK, False, FONT), ("BRUTE FORCE", 10, True, ROSE, False, MONO)]}])
txt(s, 0.94, ay + 0.74, aw - 0.6, 1.55, [
    one("遍歷：每個 entity 都要兩兩比較 → O(N²)", 12.5, True, TEXT, sa=5),
    {'runs': [("但 pattern 最多由 ", 12, False, TEXT, False, FONT),
              ("k = 120", 13, True, HEAD, False, MONO),
              (" 個 entity 組成，每個都要比 → 再 ×k", 12, False, TEXT, False, FONT)], 'space_after': 5},
    {'runs': [("→ O(N²·k) 量級 ≈ ", 12, False, TEXT, False, FONT),
              ("8.1 × 10¹³", 15, True, ROSE_DK, False, MONO),
              (" 次／模板", 12, False, TEXT, False, FONT)], 'space_after': 6},
    one("真實 DXF N ≈ 820,212；每次還要旋轉對位＋容差比對——遍歷根本跑不完。", 12, False, MUTE, ls=1.12)])

rx2 = 0.6 + aw + 0.33
card(s, rx2, ay, aw, ah, 'dark')
label(s, rx2 + 0.32, ay + 0.24, "INSPIRATION", AMBER, size=10)
txt(s, rx2 + 0.32, ay + 0.52, aw - 0.6, 0.4, [one("靈感：天文學家怎麼找星系／星團／星雲？", 13.5, True, WHITE)])
txt(s, rx2 + 0.32, ay + 1.0, aw - 0.62, 1.25, [
    one("不逐顆星互比（那是宇宙級 O(N²)）——而是：", 11.5, False, DLT, sa=5, ls=1.1),
    {'runs': [("①  ", 12, True, CYAN, False, MONO), ("旋轉／縮放／平移都不變的幾何簽章", 11.5, False, DLT, False, FONT)], 'space_after': 3},
    {'runs': [("②  ", 12, True, CYAN, False, MONO), ("空間索引快速定位候選天區", 11.5, False, DLT, False, FONT)], 'space_after': 3},
    {'runs': [("③  ", 12, True, CYAN, False, MONO), ("只對倖存候選做精確比對", 11.5, False, DLT, False, FONT)], 'space_after': 5},
    one("（同類思路：astrometry.net quad-hash、FoF 分群）", 12, False, DMUTE, ls=1.08)])

# the algorithm, applied to DXF
my = 5.06
txt(s, 0.6, my, 12.1, 0.34, [
    {'runs': [("把同一招搬到 DXF：", 13, True, HEAD, False, FONT),
              ("同形圖樣 = 一個星座，整張圖 = 一片星空", 13, True, CYAN8, False, FONT)]}])
algo = [
    ("① 幾何簽章 gate", "旋轉/鏡像/縮放不變", 'cyan'),
    ("② cKDTree 空間索引", "O(log N) 鄰域", 'navy'),
    ("③ PCA 主軸對位", "sign-variant 擇優", 'cyan'),
    ("④ 對稱 chamfer", "容差評分", 'navy'),
    ("⑤ inverted-index", "內含抑制→近線性", 'cyan'),
]
aw3 = (12.13 - 4 * 0.34) / 5
ya = my + 0.4
for i, (t, sub, kind) in enumerate(algo):
    x = 0.6 + i * (aw3 + 0.34)
    fill = CYAN6 if kind == 'cyan' else NAVY
    fg = NAVY if kind == 'cyan' else WHITE
    sub_c = RGBColor(0x0B, 0x3A, 0x44) if kind == 'cyan' else DLT
    rect(s, x, ya, aw3, 0.92, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, shadow=True)
    txt(s, x + 0.08, ya + 0.16, aw3 - 0.16, 0.42, [one(t, 11.5, True, fg, align=PP_ALIGN.CENTER, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.08, ya + 0.57, aw3 - 0.16, 0.3, [one(sub, 9.5, False, sub_c, align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    if i < 4:
        rect(s, x + aw3 + 0.06, ya + 0.33, 0.22, 0.26, fill=CYAN6, shape=MSO_SHAPE.CHEVRON)

# payoff strip — two distinct pipeline stages, each with its own N (not one curve)
card(s, 0.6, 6.46, 12.13, 0.74, 'navy')
txt(s, 0.92, 6.46, 11.5, 0.74, [
    {'runs': [("遍歷 brute-force　", 12.5, True, ROSE, False, FONT),
              ("O(N²·k) ≈ 8.1×10¹³ 次／模板（k≤120）　→　單核估算約 數週～數月", 12.5, True, WHITE, False, FONT)], 'space_after': 4},
    {'runs': [("尋形 Conform　", 12.5, True, CYAN, False, FONT),
              ("簽章 gate ＋ cKDTree 空間索引剪枝到近線性　→　互動式秒級回應", 12.5, True, WHITE, False, FONT)]}],
    anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 1:38–2:30（52 秒）——演算法 + 複雜度 + 天文學靈感頁。\n"
         "先把問題講成幾何模式比對：在一張 DXF（N≈820,212 個實體）裡，找某個框選圖樣的『全部同形實例』。\n"
         "暴力解 O(N²·k)：遍歷每個 entity 兩兩比較是 O(N²)；又因 query pattern 最多由 k=120 個 entity 組成、每個都要比 → 再 ×k → O(N²·k)。"
         "真實 DXF N≈820,212、k≤120 → 量級 ≈ 8.1×10¹³ 次／每個模板，每次還要旋轉對位＋容差比對，遍歷下完全跑不完"
         "（8.1×10¹³ 是 O(N²·k) 量級估算示意，N 取自真實 DXF、k 為目前 pattern 上限）。\n"
         "遍歷時間重新估算：單核以 ~10⁷–10⁸ 次/秒（含對位＋容差的複合比較）計，8.1×10¹³ 次 ≈ 8.1×10⁶–8.1×10⁵ 秒 ≈ 數週～數月；"
         "對照尋形 Conform 經簽章 gate + cKDTree 剪枝後為互動式秒級（吞吐假設依單次比較成本而定，為量級估算非現場 benchmark）。\n"
         "靈感（口述可展開）：天文學家找星系／星團／星雲不是逐顆星互比，而是用對旋轉/縮放/平移不變的幾何特徵『簽章』描述結構、"
         "丟進空間索引先定位候選天區、再只對倖存者做精確比對——例如 astrometry.net 用 4 星 quad 的幾何 hash 做 blind plate-solving"
         "（對 平移/旋轉/縮放/鏡射 不變），friends-of-friends 用連結長度分群找星團。\n"
         "尋形 Conform 把同一套思路搬到 DXF：同形圖樣＝一個星座，整張圖＝一片星空。\n"
         "演算法五步：①幾何簽章 gate（旋轉/鏡像/縮放不變）剪掉絕大多數候選 → ②cKDTree 空間索引 O(log N) 鄰域 → "
         "③PCA 主軸對位＋4 個鏡射/旋轉 sign-variant 取最佳（Procrustes 式對齊）→ ④對稱 chamfer 評分容差 → ⑤inverted-index 內含抑制。\n"
         "效能：比對掃描 O(N²·k) 遍歷（單核估算數週～數月）→ 經簽章 gate + cKDTree 空間索引剪枝到近線性、互動式秒級回應。"
         "（內含抑制 20s→8ms 是另一後處理步、見 P5；皆為 design.md／備忘的實測或量級估算，非現場 benchmark。）\n"
         "誠實措辭：astrometry.net / friends-of-friends 是『同類方法』的類比，用來說明思路，非本專案直接套用其程式庫；"
         "本專案實作為自寫的 signature gate + scipy cKDTree + PCA 主軸對位/sign-variant 擇優/對稱 chamfer + inverted index（matcher 為 Procrustes 式對位，未跑 SVD-based Kabsch）。")

# =================================================================
# Slide 5 — 成效與效益 (★ 40%)
# =================================================================
s = slide()
header(s, "成效 ★", "IMPACT", "成效：6 人 10 天 → 1 人 1 小時",
       "Results & Impact — the core metric", 5, accent=AMBER, badge_w=2.55)

# hero metric band (dark, reference LOSS-card spirit)
card(s, 0.6, 2.16, 12.13, 1.42, 'dark')
label(s, 0.95, 2.3, "VERIFY ONE DRAWING", CYAN, size=10)
txt(s, 0.92, 2.56, 7.15, 0.95, [
    {'runs': [("6 人 × 10 天", 24, True, ROSE, False, FONT),
              ("   →   ", 22, True, WHITE, False, MONO),
              ("1 人 × 1 小時", 26, True, EMER_LT, False, FONT)],
     'align': PP_ALIGN.LEFT, 'space_after': 3},
    one("＝ 約 480 等效人工工時（6人×10天×8h）→ 約 1 小時機器運行", 12, True, CYAN)], anchor=MSO_ANCHOR.TOP)
rect(s, 8.15, 2.42, 4.32, 0.9, fill=CYAN6, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 8.15, 2.42, 4.32, 0.9, [
    one("≈ 480 → 1", 26, True, WHITE, align=PP_ALIGN.CENTER, sa=1, name=MONO),
    one("人力投入約 1/480（工程師小時）", 11.5, False, RGBColor(0xD7, 0xF0, 0xF6), align=PP_ALIGN.CENTER),
], anchor=MSO_ANCHOR.MIDDLE)

# before/after table
ty = 3.82
col = [0.6, 3.0, 7.6, 12.73]
rows = [
    ("人力 / 工時", "6 位資深工程師 × 10 天", "1 位工程師 × 1 小時（含操作 + 100+ 條自動化條件檢查）"),
    ("覆蓋率", "抽樣（缺陷會漏）", "100% 全 entity 全覆蓋 ＋ 100+ 條自動化檢查（架構保證）"),
    ("併發 / 吞吐", "4 套 AutoCAD 授權序列化", "多圖同時處理，不被單張圖卡住；零 AutoCAD 授權"),
]
rh = 0.72
rect(s, col[0], ty, col[3] - col[0], 0.5, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, col[0] + 0.2, ty, col[1] - col[0] - 0.2, 0.5, [one("維度", 13, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
txt(s, col[1] + 0.2, ty, col[2] - col[1] - 0.2, 0.5, [one("BEFORE（舊流程）", 13, True, RGBColor(0xFF, 0xB3, 0xC2))], anchor=MSO_ANCHOR.MIDDLE)
txt(s, col[2] + 0.2, ty, col[3] - col[2] - 0.2, 0.5, [one("AFTER（尋形 Conform）", 13, True, EMER_LT)], anchor=MSO_ANCHOR.MIDDLE)
for i, (lab, bef, aft) in enumerate(rows):
    yy = ty + 0.5 + i * rh + 0.06
    rect(s, col[0], yy, col[1] - col[0], rh - 0.04, fill=GRAYBG, line=CARDLN, line_w=0.75)
    rect(s, col[1], yy, col[2] - col[1], rh - 0.04, fill=ROSE_LT, line=CARDLN, line_w=0.75)
    rect(s, col[2], yy, col[3] - col[2], rh - 0.04, fill=EMER_LT, line=CARDLN, line_w=0.75)
    txt(s, col[0] + 0.2, yy, col[1] - col[0] - 0.25, rh - 0.04, [one(lab, 13, True, NAVY)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, col[1] + 0.2, yy, col[2] - col[1] - 0.3, rh - 0.04, [one(bef, 12.5, False, ROSE_DK)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, col[2] + 0.2, yy, col[3] - col[2] - 0.3, rh - 0.04, [one(aft, 12.5, True, EMER)], anchor=MSO_ANCHOR.MIDDLE)

# headline + scale strip
hy = ty + 0.5 + 3 * rh + 0.18
rect(s, 0.6, hy, 7.5, 0.72, fill=AMBER_LT, line=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.85, hy, 7.1, 0.72, [one("「抽樣 → 選定圖樣 100% 全掃」是品質階躍——漏檢從機率，變成可量測的覆蓋率。", 13, True, RGBColor(0x8A, 0x59, 0x06))],
    anchor=MSO_ANCHOR.MIDDLE)
card(s, 8.3, hy, 4.43, 0.72, 'navy', radius=0.12)
txt(s, 8.55, hy, 4.0, 0.72, [
    one("生產級規模｜真實 DXF 820,212 實體", 12, True, CYAN, sa=1),
    one("內含抑制 20s → 8ms（實測）", 11.5, True, DLT)], anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 2:30–3:38（68 秒）——★ 最高權重(40%)的核心頁，講最久。\n"
         "頭條：6 人 × 10 天 → 1 人 × 1 小時（含操作 + 100+ 條自動化檢查）。換算約 480 工程師小時 → 1，約 1/480。\n"
         "三列對比表逐列講：(1)人力工時；(2)覆蓋率——這是最防守得住的一句：從抽樣變 100% 全 entity，"
         "再加 100+ 條自動化條件檢查，『漏檢』從機率問題變成不可能，這是品質的階躍不只是加速；"
         "(3)併發/吞吐——多圖同時處理，不被單一張圖紙的檢驗卡住，且不再受 4 套 AutoCAD 授權序列化。\n"
         "規模佐證(實測)：真實 DXF 含 820,212 個實體；內含抑制 O(N²)→倒排索引，20s→8ms。\n"
         "誠實守則：1 小時是實際運行成效；480× 是換算，被問就說明算法(6×10×8 工時 vs 1 工時)。")

# =================================================================
# Slide 6 — 核心理念：LLM 時代怎麼做出「對」的產品
# =================================================================
s = slide()
header(s, "理念", "THESIS", "LLM 時代，怎麼做出「對」的產品？",
       "Anyone can ship with an LLM — shipping the right thing is the hard part", 6,
       accent=VIOLET, badge_w=2.3)

txt(s, 0.6, 2.16, 12.1, 0.5, [
    one("LLM 把「寫程式」變便宜了，卻沒有把「做對產品」變便宜——需求、品質、AI 的發揮，仍然取決於工程紀律與工具。",
        14.5, True, TEXT)])

thesis = [
    ("真的符合需求嗎？", "規格先行", CYAN6,
     ["意圖先寫成 proposal / design，在紙上吵完再動手——", "方向錯了在規格層修，成本是幾行字；", "等 diff 出來才修，成本是整個功能重寫。"]),
    ("品質怎麼確保？", "流程強制", EMER,
     ["品質不是宣稱，是閘門——機器驗證的能力契約、", "605 個測試、真瀏覽器端到端驗證，", "每一關不過就不能 merge。"]),
    ("怎麼讓 AI 表現更好？", "給對工具", AMBER,
     ["codegraph 給它導航、playwright 給它眼睛、", "skills 把流程固化成生產線——", "AI 負責提案，工程判斷負責把關。"]),
]
tw6 = (12.13 - 2 * 0.33) / 3
ty6 = 2.78; th6 = 3.3
for i, (q, kw, c, body) in enumerate(thesis):
    x = 0.6 + i * (tw6 + 0.33)
    card(s, x, ty6, tw6, th6, 'white')
    rect(s, x, ty6 + 0.18, 0.12, th6 - 0.36, fill=c)
    txt(s, x + 0.32, ty6 + 0.3, tw6 - 0.5, 0.5, [one(f"Q{i+1}　{q}", 15, True, HEAD)])
    chip(s, x + 0.32, ty6 + 0.92, 1.6, 0.42, kw, c, WHITE if c != AMBER else NAVY, size=13)
    txt(s, x + 0.32, ty6 + 1.56, tw6 - 0.6, th6 - 1.8,
        [one(t, 12, False, TEXT, ls=1.15, sa=2) for t in body])

card(s, 0.6, 6.4, 12.13, 0.78, 'navy')
txt(s, 0.95, 6.4, 11.4, 0.78, [
    one("這三個問題，就是接下來兩頁——尋形 Conform 的開發生產線與品質工程，全部圍繞它們設計。",
        14.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 3:38–4:08（30 秒）——核心理念頁，給後半場定調。\n"
         "一句開場：LLM 時代 demo 很便宜，人人都能『做出一個產品』；但『能跑』不等於『對』。"
         "寫程式的成本降了，做對產品的成本沒降。\n"
         "三問三答各一句：(Q1) 符合需求嗎→規格先行，在 proposal/design 紙上吵完才動手；"
         "(Q2) 品質怎麼確保→流程強制，契約+測試+真瀏覽器驗證當閘門，不靠宣稱；"
         "(Q3) 怎麼讓 AI 更好→給對工具，codegraph 導航/playwright 眼睛/skills 生產線，AI 提案、人把關。\n"
         "收尾：接下來兩頁就是這三題在本專案的具體實作——轉 P7。")

# =================================================================
# Slide 7 — Claude Code 生產線
# =================================================================
s = slide()
header(s, "開發", "PROCESS", "用 Claude Code 把開發變成生產線",
       "Building Conform as a Claude Code Production Line", 7, accent=VIOLET, badge_w=2.3)

steps4 = [
    ("① 人類意圖\nspectra-propose", False),
    ("② codegraph 導航\nsearch · impact", False),
    ("③ Claude 實作\n編碼 + 自寫測試", False),
    ("④ playwright\n開真瀏覽器截圖驗證", False),
    ("⑤ [USER] 機密 DXF\n親驗｜信任邊界", True),
    ("⑥ spectra-archive\n歸檔 + spec 同步", False),
]
bx, by = 0.6, 2.16
bw, bh = 3.5, 0.78
hg, vg = 0.55, 0.3
pos = {
    0: (bx + 0 * (bw + hg), by),
    1: (bx + 1 * (bw + hg), by),
    2: (bx + 2 * (bw + hg), by),
    3: (bx + 2 * (bw + hg), by + bh + vg),
    4: (bx + 1 * (bw + hg), by + bh + vg),
    5: (bx + 0 * (bw + hg), by + bh + vg),
}
for i, (lab, hot) in enumerate(steps4):
    x, y = pos[i]
    if hot:
        rect(s, x, y, bw, bh, fill=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1, shadow=True)
        fg = NAVY
    else:
        card(s, x, y, bw, bh, 'tint', radius=0.1)
        fg = NAVY
    txt(s, x + 0.12, y, bw - 0.24, bh, [one(lab, 12.5, True, fg, align=PP_ALIGN.CENTER, ls=1.05)], anchor=MSO_ANCHOR.MIDDLE)
for a in (0, 1):
    x, y = pos[a]
    rect(s, x + bw + 0.04, y + bh / 2 - 0.13, hg - 0.08, 0.26, fill=VIOLET, shape=MSO_SHAPE.RIGHT_ARROW)
x, y = pos[2]
rect(s, x + bw / 2 - 0.13, y + bh + 0.02, 0.26, vg - 0.04, fill=VIOLET, shape=MSO_SHAPE.DOWN_ARROW)
for a in (3, 4):
    x, y = pos[a]
    rect(s, x - hg + 0.04, y + bh / 2 - 0.13, hg - 0.08, 0.26, fill=VIOLET, shape=MSO_SHAPE.LEFT_ARROW)
x, y = pos[5]
rect(s, x + bw / 2 - 0.13, y - vg + 0.02, 0.26, vg - 0.04, fill=VIOLET, shape=MSO_SHAPE.UP_ARROW)

# tool-stack chips (row 1) + skill ecosystem band (row 2)
tcy = by + 2 * bh + vg + 0.14
chip(s, 0.6, tcy, 3.95, 0.42, "MCP：codegraph · playwright · discord", NAVY, WHITE, size=11)
chip(s, 4.75, tcy, 3.9, 0.42, "Hooks：RTK 省 ~52% token · auto-sync", RGBColor(0x3A, 0x4A, 0x6A), WHITE, size=11)
chip(s, 8.85, tcy, 3.88, 0.42, "信任邊界：機密 DXF 不進 context", AMBER, NAVY, size=11)
sky = tcy + 0.52
card(s, 0.6, sky, 12.13, 0.62, 'tint', radius=0.1)
txt(s, 0.85, sky, 11.65, 0.62, [
    {'runs': [("Skills｜自建技能生態　", 12, True, CYAN8, False, FONT),
              ("規格流 spectra-* ×12 · 測試 tdd · 審查 grill-me / diagnose / triage",
               12, False, TEXT, False, FONT)], 'space_after': 2},
    one("規劃 to-prd / to-issues · 領域 add-rule · 雛型 prototype · 架構 improve-codebase-architecture 等",
        12, False, MUTE)], anchor=MSO_ANCHOR.MIDDLE)

# decision strip
dy = sky + 0.72
card(s, 0.6, dy, 12.13, 1.34, 'white')
txt(s, 0.85, dy + 0.1, 11.6, 0.4, [one("AI 提案，工程判斷把關——每個採用都附一個被我否決的替代方案", 13, True, NAVY)])
deci = [
    ("採用", EMER, "我否決 Claude 的「密度啟發式」，改採確定性互斥視圖規則——可重現、operator 看得懂"),
    ("後來推翻", ROSE, "密度仲裁上線後產線 17,482 球整批誤判，我判定不追不可重現的 bug、整個子系統當死碼砍掉"),
    ("被產線打臉再修", AMBER, "design 判 O(M²)「便宜」被 20k 球產線打臉到 20s+，我換 inverted index 近線性（6000 球 2.0s→8ms）"),
    ("延後", MUTE, "Postgres／MinIO 遷移我寫好 ADR 但標「未實作」——_jobs 才是硬牆，無 HA 需求前不動"),
]
for i, (tag, c, body) in enumerate(deci):
    yy = dy + 0.52 + i * 0.205
    chip(s, 0.85, yy, 1.55, 0.2, tag, c, WHITE, size=10.5)
    txt(s, 2.5, yy - 0.02, 10.1, 0.26, [one(body, 11, False, TEXT)], anchor=MSO_ANCHOR.MIDDLE)

txt(s, 0.6, dy + 1.42, 12.13, 0.34, [one(
    "55 archived + 34 active changes　·　10 capability specs　·　290 commits　·　兩個月（2026-05-13 起）",
    12, True, CYAN8, align=PP_ALIGN.CENTER, name=MONO)])

notes(s,
    "⏱ 4:08–4:55（47 秒）\n"
    "口頭講稿(約180字)：我把 Claude 當一位資深工程師，每個提案都拿來質問，不照單全收。\n"
    "舉真實功能『量測距離工具』：先下 spectra-propose，Claude 產出 proposal、design(10 條編號決策)、tasks 三份文件，"
    "我在規格層就審完意圖、便宜地修正方向，而不是事後改 diff。接著 spectra-apply，Claude 用 codegraph 直接跳到 "
    "canvas.js 節點、查呼叫關係與改動影響面，不在三千八百行裡盲目 grep；邊實作邊勾 task、自寫測試。前端是 Canvas2D，"
    "DOM 看不到像素，所以用 playwright 開真瀏覽器框選、截圖比對 overlay。最後一格永遠標 [USER]——機密客戶 DXF 只由我親驗，"
    "Claude 碰不到。驗完下 spectra-archive 自動歸檔並同步規格。兩個月、55 個歸檔功能（34 個進行中），每個都走同一條紀律化迴圈。\n\n"
    "── 完整採用/否決決策表(評審深問時口述)──\n"
    "1) 同半徑 BGA vs 對位點：Claude 提密度仲裁→我採用→產線 17,482 球誤判、合成案重現不出→改互斥視圖規則(commit 77e8832)。\n"
    "2) 失效仲裁子系統：曾主張保留當骨架→我判 YAGNI 死碼，直接刪(fffd30e)。\n"
    "3) 含括抑制：design 判 O(M²)便宜→20k 球實測 20s+→換 inverted handle-index，6000 球 2.0s→8ms(13f179f)。\n"
    "4) multi-dxf：先引 dxf_view enum(D1-7)→Decision 8 自我簡化推翻，enum 是 over-fit。\n"
    "5) 基板比對：per-class tolerance 先 revert(9db06a6→c4df21d)→改 per-class match_strategy。\n"
    "6) SEC-001 上傳預檢：延後(multipart Content-Length 涵蓋整 body，對內部信任使用者邊際)。\n"
    "7) 儲存遷移 Postgres/MinIO：ADR 標『討論中、尚未實作』，_jobs 才是多 replica 硬牆。\n\n"
    "── 誠實措辭 ──\n"
    "• codegraph 是鐵證(.codegraph/codegraph.db 5.5MB 實建、7 工具預授權)。\n"
    "• playwright/discord 說『啟用於我的 Claude Code 環境並用於本專案』，不要說『專案檔裡配置』(專案無 .mcp.json)。\n"
    "• playwright 視覺回歸、discord scan-all 通知是『使用方式』描述，非既存測試套件；被追問誠實說明。\n"
    "• 絕不可宣稱 Gmail/Calendar/Drive(未接入，config 中反列 OAuth-blocked)。真正三件套：codegraph + playwright + discord。\n"
    "• RTK 省 ~52% token(4.8M/88k 指令)是全域 hook 分析數據。")

# =================================================================
# Slide 8 — 程式品質 × 規格驅動
# =================================================================
s = slide()
header(s, "品質", "QUALITY", "程式品質 × 規格驅動工程",
       "Code Quality × Spec-Driven Development", 8, accent=CYAN8, badge_w=2.3)

colw = 5.9
lx, ly = 0.6, 2.18
card(s, lx, ly, colw, 2.95, 'tint')
txt(s, lx + 0.3, ly + 0.2, colw - 0.6, 0.5, [one("規格驅動：品質由流程強制，非口頭宣稱", 15, True, CYAN8)])
for i, t in enumerate([
    "10 個機器驗證能力契約（SHALL + WHEN/THEN 場景）",
    "89 個 OpenSpec changes（55 archived + 34 active）",
    "openspec validate --specs 為 pre-merge 閘門",
    "每個 change 一分支、留書面決策與審計軌跡",
]):
    txt(s, lx + 0.3, ly + 0.82 + i * 0.5, colw - 0.6, 0.45, [
        {'runs': [("▸ ", 13.5, True, CYAN6, False, FONT), (t, 13, False, TEXT, False, FONT)]}])

rx = lx + colw + 0.33
card(s, rx, ly, colw, 2.95, 'white')
txt(s, rx + 0.3, ly + 0.2, colw - 0.6, 0.5, [one("工程硬證據（每條釘 file / test）", 15, True, NAVY)])
for i, t in enumerate([
    "605 測試函式、測試 13.3k LOC > 應用 12.9k LOC",
    "倒排索引：含括抑制 O(N²) → 近線性（守 N=20000 <2s）",
    "AST 不變量測試：workers 永不讀 stale cache",
    "DRC envelope 嚴格驗證：壞輸出 → 大聲 400，不靜默",
    "canvas.js ↔ Python 常數 drift-guard 測試",
]):
    txt(s, rx + 0.3, ly + 0.82 + i * 0.41, colw - 0.6, 0.4, [
        {'runs': [("✓ ", 13, True, EMER, False, FONT), (t, 12.5, False, TEXT, False, FONT)]}])

# matcher band (dark)
my = 5.28
card(s, 0.6, my, 12.13, 1.0, 'dark')
label(s, 0.92, my + 0.14, "MATCHER", CYAN, size=10)
txt(s, 2.3, my + 0.12, 10.2, 0.85, [
    one("PCA 主軸對位 + sign-variant 擇優 + 對稱 chamfer → 平移／旋轉／鏡像／±5% 縮放不變", 13.5, True, WHITE, sa=2),
    one("三條「真實 bug 驅動」fast path（820,212 實體 DXF bucket-split、O(N²) 抑制爆 20s）——bug 驅動的演算法，不是紙上設計。",
        12, False, DLT)])

# callback line
rect(s, 0.6, 6.42, 12.13, 0.66, fill=AMBER_LT, line=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 6.42, 11.5, 0.66, [one(
    "規格驅動讓「不同意」變便宜——我們在 design.md 的 ## Decisions 吵完每個 Alternative，而不是等 diff 出來才回頭。",
    13, True, RGBColor(0x8A, 0x59, 0x06))], anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 4:55–5:38（43 秒）——把 25%+25% 兩個品質軸(程式品質 + 開發流程)合在這頁高密度交代。\n"
         "左欄：規格驅動——品質是被流程『強制』出來的(10 契約、89 changes、validate 閘門)，不是嘴上說有測試。\n"
         "右欄：工程硬證據，每條都能翻到對應 file/test(倒排索引、AST 不變量、envelope 驗證、drift-guard)。\n"
         "中段藍帶：matcher 是技術核心，三條 fast path 由真實生產 bug 逼出(820,212 實體)——強調『bug 驅動，不是紙上設計』。\n"
         "結尾 callback 呼應 P7 的決策表：spec-driven 讓爭論發生在 design.md，不是 diff 之後。\n"
         "誠實：605 是 tests/ 下 test 函式數(2026-07-14 重新統計；含參數化展開更多)；前端 canvas 零自動化 UI 測試，被問就框成 roadmap。"
         "呼應 P6 三問：左欄=Q1 規格先行、右欄=Q2 流程強制。")

# =================================================================
# Slide 9 — 創新差異化與未來擴展（收尾）
# =================================================================
s = slide()
header(s, "擴展", "ROADMAP", "創新差異化與未來擴展",
       "Differentiation & Roadmap", 9, accent=EMER, badge_w=2.3)

# vs old method
card(s, 0.6, 2.18, 12.13, 0.9, 'white')
txt(s, 0.9, 2.18, 11.5, 0.9, [
    {'runs': [("vs 舊做法　", 14, True, ROSE, False, FONT),
              ("人工逐物件抽樣量測（AutoCAD 內手動）", 13, False, TEXT, False, FONT),
              ("　→　", 14, True, NAVY, False, FONT),
              ("框選 → 比對 → 規則檢查的閉環（class-agnostic 幾何引擎，非封裝寫死）", 13, True, EMER, False, FONT)],
     'line_spacing': 1.1}], anchor=MSO_ANCHOR.MIDDLE)

# radial expansion
cx_, cy_ = 3.4, 4.72
spx = 5.55; spw = 7.18; sph = 0.56; pitch = 0.62; sy0 = 3.22
spokes = [
    ("先進封裝 DXF", "已實作", EMER, EMER_LT),
    ("PCB Gerber / ODB++　間距 DRC", "可擴展", CYAN6, TINT),
    ("GDSII / OASIS　layout 圖樣列舉", "可擴展", CYAN6, TINT),
    ("photomask / MEMS　圖樣稽核", "可擴展", CYAN6, TINT),
    ("interposer / RDL　軟板佈局檢查", "可擴展", CYAN6, TINT),
]
for i in range(len(spokes)):
    yy = sy0 + i * pitch
    cmid = yy + sph / 2
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cx_ + 1.2), Inches(cy_),
                                Inches(spx), Inches(cmid))
    cn.line.color.rgb = CARDLN; cn.line.width = Pt(1.5)
# center hub
rect(s, cx_ - 1.25, cy_ - 0.72, 2.5, 1.44, fill=NAVY, shape=MSO_SHAPE.OVAL, shadow=True)
brandmark(s, cx_ - 0.27, cy_ - 0.52, 0.5, form_color=CYAN, bracket_color=WHITE)
txt(s, cx_ - 1.25, cy_ + 0.02, 2.5, 0.66, [
    one("尋形 Conform", 16, True, WHITE, align=PP_ALIGN.CENTER, sa=1),
    one("幾何 → 規則引擎", 11.5, False, DMUTE, align=PP_ALIGN.CENTER)],
    anchor=MSO_ANCHOR.MIDDLE)
for i, (name, tag, c, lt) in enumerate(spokes):
    yy = sy0 + i * pitch
    rect(s, spx, yy, spw, sph, fill=lt, line=c, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16, shadow=True)
    txt(s, spx + 0.28, yy, spw - 1.55, sph, [one(name, 12.5, True, TEXT, ls=1.0)],
        anchor=MSO_ANCHOR.MIDDLE)
    chip(s, spx + spw - 1.3, yy + (sph - 0.34) / 2, 1.05, 0.34, tag, c, WHITE, size=11)

# closing band — production hardening in flight + positioning
card(s, 0.6, 6.34, 12.13, 0.9, 'navy')
txt(s, 0.9, 6.34, 11.5, 0.9, [
    {'runs': [("生產化進行中　", 12.5, True, AMBER, False, FONT),
              ("MariaDB 多副本 · MinIO 物件儲存 · Keycloak SSO 三級權限 · 產品版本化 · 簽核存證（openspec 34 active changes）",
               12.5, True, WHITE, False, FONT)], 'space_after': 4},
    {'runs': [("通用定位　", 12.5, True, CYAN, False, FONT),
              ("DRC 引擎為 adapter + 版本化 bundle → 可獨立出貨；「幾何→規則引擎」通用前端層，可餵任意 DRC 後端。",
               12.5, True, WHITE, False, FONT)]}],
    anchor=MSO_ANCHOR.MIDDLE)

notes(s, "⏱ 5:38–6:00（22 秒）——創新 10% 核心 + 應用廣度收尾。\n"
         "vs 舊法：一句話講清差異——從人工逐物件抽樣量測，變成框選→比對→規則檢查的閉環，引擎是 class-agnostic 幾何不是封裝寫死。\n"
         "輻射圖：中央是 尋形 Conform『幾何→規則引擎』，向外是鄰近應用——封裝 DXF 已實作，其餘(PCB Gerber/ODB++、GDSII/OASIS、"
         "photomask/MEMS、interposer/RDL)標『可擴展』，誠實不誇大成已做。\n"
         "底帶第一行是『進行中』的生產化硬化（openspec add-production-infra-and-auth / add-product-versioning / "
         "add-signoff-evidence 等 34 個 active changes）——MariaDB 多副本、MinIO、Keycloak SSO 三級權限(admin/editor/viewer × "
         "global/customer/product)、一版一庫的產品版本化、簽核存證。措辭用『進行中』，不可說已上線。\n"
         "底帶第二行：乾淨邊界(adapter + 版本化 bundle)讓 尋形 Conform 可獨立出貨，定位成多個 DRC 後端的通用前端層。\n"
         "收尾口頭金句可用：『尋形 Conform 不是一個封裝檢查工具，是一層幾何→規則的通用前端。』")

# =================================================================
# Slide 9 — Q&A（3 分鐘）
# =================================================================
s = slide()
rect(s, 0, 0, SW, SH, fill=BG_DARK)
rect(s, 0, SH - 0.16, SW, 0.16, fill=CYAN6)
brandmark(s, 6.17, 1.3, 1.0, form_color=CYAN, bracket_color=WHITE)
txt(s, 0.6, 2.55, SW - 1.2, 1.0, [one("Q & A", 44, True, WHITE, align=PP_ALIGN.CENTER, name=MONO)])
txt(s, 0.6, 3.6, SW - 1.2, 0.5, [
    one("6 人 × 10 天 → 1 人 × 1 小時　·　選定圖樣 100% 覆蓋　·　零 AutoCAD 授權依賴", 16, True, DLT, align=PP_ALIGN.CENTER)])
qa_hints = [("演算法細節", "簽章 gate / cKDTree / PCA 對位 → P4"),
            ("數字怎麼算", "480 工時換算與量級估算 → P5 備註"),
            ("AI 怎麼用", "採用／否決決策日誌 → 下一頁 BACKUP")]
qw = 3.6; qg = 0.45
qx0 = (SW - 3 * qw - 2 * qg) / 2
for i, (t, sub) in enumerate(qa_hints):
    x = qx0 + i * (qw + qg)
    rect(s, x, 4.55, qw, 1.15, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12, shadow=True)
    txt(s, x, 4.73, qw, 0.45, [one(t, 15, True, CYAN, align=PP_ALIGN.CENTER)])
    txt(s, x + 0.2, 5.18, qw - 0.4, 0.4, [one(sub, 10.5, False, DMUTE, align=PP_ALIGN.CENTER)])
txt(s, SW - 3.5, 0.35, 3.0, 0.3, [one("Q&A · 3 MIN", 11, True, CYAN, align=PP_ALIGN.RIGHT, name=MONO)])
notes(s, "⏱ 6:00 後 · Q&A（另計，約 3 分鐘）\n"
         "停在這頁等提問；一行 recap 讓評審提問時眼前就有核心數字。\n"
         "預判三類問題與跳頁策略：演算法細節→回 P4；數字換算→口述 P5 備註的算法(6×10×8 工時 vs 1 工時、量級估算非現場 benchmark)；"
         "AI 使用深度→切下一頁 BACKUP 決策日誌逐列講。\n"
         "誠實守則：效能絕對數字以 design.md／專案備忘為準；被問到沒把握的數字，答『量級估算』不要硬報精確值。")

# =================================================================
# Slide 10 — BACKUP：完整採用／否決決策日誌（Q&A 深問時用）
# =================================================================
s = slide()
header(s, "附錄", "APPENDIX", "採用／否決 決策日誌",
       "AI-proposed / human-decided log（Q&A backup）", 0, accent=MUTE, badge_w=2.3, tag="BACKUP")

cx = [0.5, 2.5, 4.7, 6.12, 10.45, 12.83]
heads = ["改進主題", "Claude 建議", "我的決定", "我們怎麼討論／理由", "結果"]
ty = 2.12; hh = 0.46
rect(s, cx[0], ty, cx[5] - cx[0], hh, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
for i, h in enumerate(heads):
    txt(s, cx[i] + 0.12, ty, cx[i + 1] - cx[i] - 0.18, hh, [one(h, 12.5, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

DEC = [
    ("同半徑 BGA 球 vs 對位點", "鄰域密度仲裁（自動推 pitch、鄰居數分類）", "後來推翻", ROSE,
     "初期採用（無參數）；產線 17,482 球大陣列整批誤判成 Fiducial、合成案重現不出 → 不追不可重現的啟發式 bug", "改用視圖約束"),
    ("替代分類規則", "互斥視圖判定（BGA=底視圖、Fiducial=頂視圖）", "採用", EMER,
     "確定性、operator 看得懂、免疫密度 edge case；明確否決「保留密度群當保險」", "上線穩定（77e8832）"),
    ("失效仲裁子系統去留", "保留 arbitrate() 當未來骨架", "改寫（YAGNI 刪）", AMBER,
     "視圖約束讓 registry 清空、call site 全成 pass-through → 可達但無作用的死碼，直接刪", "子系統移除（fffd30e）"),
    ("含括重複比對抑制效能", "每類 pairwise 掃描，design 判 O(M²)「便宜」", "採用後修正", AMBER,
     "上線打臉：20k 球單這步 20s+ → 換 inverted handle-index（X 的超集必含 X 全部 handle），近線性", "6000 球 2.0s→8ms（13f179f）"),
    ("多 DXF 一角色的 view 維度", "引入 dxf_view enum，正交於設計角色（Decision 1-7）", "後來推翻", ROSE,
     "Decision 8 簡化轉向：使用者回饋主流程是一檔含 top/bottom/side region，分割軸未必對齊 → enum 是 over-fit", "全檔 multi，欄位 vestigial"),
    ("基板比對策略", "新增 per-class match_strategy（chamfer/signature）", "採用（取代被否決路徑）", EMER,
     "前案 per-class tolerance 已 revert（9db06a6→c4df21d）：tolerance 橋不過固定 ±20% pre-filter", "match_strategy 正式取代"),
    ("上傳上限 SEC-001 預檢", "加請求層 Content-Length 預檢，buffer 前拒絕", "延後", MUTE,
     "multipart 的 Content-Length 涵蓋整個 body 非單檔，FastAPI 無 per-part length；對內部信任使用者邊際價值低", "緩衝後 len() 為權威防線"),
]
rh = 0.60
y = ty + hh + 0.05
for r, (topic, prop, dec, dcol, why, res) in enumerate(DEC):
    bg = WHITE if r % 2 == 0 else GRAYBG
    rect(s, cx[0], y, cx[5] - cx[0], rh, fill=bg, line=CARDLN, line_w=0.5)
    txt(s, cx[0] + 0.12, y, cx[1] - cx[0] - 0.2, rh, [one(topic, 11.5, True, NAVY, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx[1] + 0.12, y, cx[2] - cx[1] - 0.2, rh, [one(prop, 10.5, False, TEXT, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx[2] + 0.1, y, cx[3] - cx[2] - 0.14, rh, [one(dec, 11, True, dcol, align=PP_ALIGN.CENTER, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx[3] + 0.12, y, cx[4] - cx[3] - 0.2, rh, [one(why, 10.5, False, TEXT, ls=1.02)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx[4] + 0.12, y, cx[5] - cx[4] - 0.2, rh, [one(res, 10.5, True, CYAN8, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    y += rh

rect(s, cx[0], y + 0.1, cx[5] - cx[0], 0.5, fill=AMBER_LT, line=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, cx[0] + 0.2, y + 0.1, cx[5] - cx[0] - 0.4, 0.5, [one(
    "主軸：「AI 提案，工程判斷把關——每個採用都附一個被我否決的替代方案」。"
    "錨點均經 git/reflog 核對；效能絕對數字以 design.md／專案備忘為準，非現場 benchmark。",
    11, True, RGBColor(0x8A, 0x59, 0x06))], anchor=MSO_ANCHOR.MIDDLE)

notes(s, "BACKUP 頁——不在 6 分鐘主線內，只在 Q&A 被深問『哪些採用哪些否決、怎麼討論』時切過來逐列講。\n"
         "P7 決策帶只放 4 條精簡版；這張是完整 7 列表。\n"
         "字級：全簡報由 _setfont 強制 ≥12pt（MIN_PT），此表理由欄 12pt 兩行內收納。\n"
         "誠實：commit hash(77e8832/fffd30e/13f179f/9db06a6→c4df21d)、Decision 8、SEC-001 皆 repo 實證；"
         "17,482 / 6000 / 20000 球與時間數字來自 design.md／專案備忘的實測記述，非現場重跑。")

prs.save("/Users/yi-tingli/Documents/Projects/side/SMDR2/competition/尋形Conform_競賽簡報.pptx")
print("saved:", len(prs.slides._sldIdLst), "slides (封面+8內容+Q&A+BACKUP)")
print("OK")
